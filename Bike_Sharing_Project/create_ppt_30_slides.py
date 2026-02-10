from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Create a presentation object
prs = Presentation()

def add_full_slide(title, points, image=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    # Body text
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.word_wrap = True
    for p in points:
        p_obj = tf.add_paragraph()
        p_obj.text = p
        p_obj.level = 0
        p_obj.space_before = Pt(10)

    if image and os.path.exists(image):
        body_shape.width = Inches(4.5)
        slide.shapes.add_picture(image, Inches(5), Inches(1.5), height=Inches(4.5))

# --- SLIDE 1: Title ---
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Bike-sharing Rental Demand Prediction"
slide.placeholders[1].text = "Data Science Project Report"

# --- SLIDE 2: Team ---
add_full_slide("Team Information", ["Mentor: [Mentor Name]", "Team Member 1: [Name]", "Team Member 2: [Name]"])

# --- SLIDE 3: Objective ---
add_full_slide("Objective", ["Predict hourly bike rental demand.", "Optimize urban city bike supply.", "Reduce customer waiting time."])

# --- SLIDE 4: Dataset Details ---
add_full_slide("Dataset Details", ["Source: UCI Bike Sharing Dataset", "17,379 Hourly Records", "17 Initial Features"])

# --- SLIDE 5: Column Descriptions 1 ---
add_full_slide("Data Dictionary: Time", ["instant: Index", "dteday: Date", "season: Seasonality", "yr: Year", "mnth: Month"])

# --- SLIDE 6: Column Descriptions 2 ---
add_full_slide("Data Dictionary: Weather", ["weathersit: Situation", "temp: Temperature", "hum: Humidity", "windspeed: Speed"])

# --- SLIDE 7: EDA - Missing Values ---
add_full_slide("EDA: Missing Value Detection", ["Detected '?' in weather metrics", "Imputed 11 entries in temp", "Imputed 6 in humidity/atemp"])

# --- SLIDE 8: EDA - Outliers ---
add_full_slide("EDA: Outlier Detection", ["Boxplots revealed high count outliers", "Confirmed as rush-hour peaks", "Kept outliers to maintain real-world signal"])

# --- SLIDE 9: EDA - Cleaning Procedure ---
add_full_slide("Data Cleaning Summary", ["Coerced numeric types", "Parsed DD-MM-YYYY dates", "Handled inconsistencies between casual/registered/cnt"])

# --- SLIDE 10: Viz - Rental Dist ---
add_full_slide("Inference: Demand Distribution", ["Right-skewed target variable", "Majority of hours are low-demand", "High peaks are rare but critical"], "images/dist_cnt.png")

# --- SLIDE 11: Viz - Hourly Trend ---
add_full_slide("Inference: Hourly Trends", ["Peaks at 8 AM and 5-6 PM", "Commuter-driven behavior dominates", "Operations should focus on these slots"], "images/hourly_trend.png")

# --- SLIDE 12: Viz - Seasonal ---
add_full_slide("Inference: Seasonal Impact", ["Fall shows highest bike demand", "Spring has significantly lower demand", "Marketing should target transition seasons"], "images/seasonal_demand.png")

# --- SLIDE 13: Viz - Weather ---
add_full_slide("Inference: Weather Sit", ["Clear weather = Maximum rentals", "Heavy rain = Minimal rentals", "Demand is highly weather-elastic"], "images/weather_impact.png")

# --- SLIDE 14: Viz - Temp ---
add_full_slide("Inference: Temperature Influence", ["Positive correlation with rentals", "Optimal riding temp: 0.6-0.8 range", "Cold weather reduces casual ridership"], "images/temp_vs_count.png")

# --- SLIDE 15: Viz - Workingday ---
add_full_slide("Inference: Working Days", ["Higher consistent demand on workdays", "Different patterns for holidays vs workdays", "Resource allocation must be dynamic"], "images/workingday_holiday.png")

# --- SLIDE 16: Feature Engineering 1 ---
add_full_slide("Feature Engineering: Cyclic", ["Hour mapped to Sin/Cos", "Month mapped to Sin/Cos", "Maintains circular time relationships"])

# --- SLIDE 17: Feature Engineering 2 ---
add_full_slide("Feature Engineering: Encoding", ["One-Hot Encoding for Season", "One-Hot Encoding for Weather", "Ensured non-ordinal treatment"])

# --- SLIDE 18: Feature Engineering 3 ---
add_full_slide("Feature Engineering: Scaling", ["Min-Max Scaler for weather data", "Normalized range: 0 to 1", "Ensured uniform feature contribution"])

# --- SLIDE 19: Correlation Matrix ---
add_full_slide("Correlation Insights", ["Temp & Atemp are highly correlated", "Humidity negatively impacts demand", "Hour is the strongest temporal feature"], "images/correlation_matrix_processed.png")

# --- SLIDE 20: Model Building Approach ---
add_full_slide("Model Selection Strategy", [
    "Evaluated three core algorithms: Decision Tree, Random Forest, and Gradient Boosting.",
    "Goal: Identify the model with the highest generalization power.",
    "Used a 80/20 train-test split with random state for reproducibility.",
    "Metric Focus: R-squared (Accuracy) and MAE (Average Error)."
])

# --- SLIDE 21: Why Not Decision Tree/Gradient Boosting? ---
add_full_slide("Analysis: Why Others Were Rejected", [
    "Decision Tree (R2: 0.88): Overfitted the training data too easily. High variance led to unreliable predictions for unseen peaks.",
    "Gradient Boosting (R2: 0.88): Showed high sensitivity to noise in the dataset, resulting in the highest Average Error (MAE 42).",
    "Both models failed to capture the complex, non-linear interactions as effectively as the ensemble method."
], "images/model_comparison_r2.png")

# --- SLIDE 22: Final Choice: Random Forest ---
add_full_slide("The Winner: Random Forest (R2: 0.94)", [
    "Ensemble Advantage: Combines multiple trees to cancel out individual prediction errors (Bagging).",
    "Robustness: Handles skewed distributions and outliers (peak hours) without losing accuracy.",
    "Performance: Achieved the lowest error and highest precision in demand forecasting."
])

# --- SLIDE 23: Hyperparameter Tuning ---
add_full_slide("Hyperparameter Tuning", ["Used Randomized Search CV", "Optimized Depth and Estimators", "Final model is robust to overfitting"])

# --- SLIDE 24: Feature Importance ---
add_full_slide("Feature Importance", ["Hour is the top driver (60%+)", "Temp is the second major driver", "Yearly growth is significant"], "images/top_drivers.png")

# --- SLIDE 25: Evaluation metrics ---
add_full_slide("Final Model Evaluation", ["MAE: ~25 rentals", "RMSE: ~43 rentals", "R-squared: 0.94 (Excellent fit)"])

# --- SLIDE 26: Actual vs Predicted ---
add_full_slide("Prediction Performance", ["Predictions follow actual values closely", "Stable across all demand levels", "Minimal error for peak hours"], "images/actual_vs_predicted.png")

# --- SLIDE 27: APP Screenshot ---
add_full_slide("Deployment: App Page", ["Streamlit Interface integrated", "Real-time user inputs supported", "Instant demand forecasting rendered"])

# --- SLIDE 28: Challenges Faced ---
add_full_slide("Challenges Faced", ["Data cleaning of hidden '?' values", "Capturing cyclic nature of time", "Tuning ensemble weights", "Managing high-dimensional categorical data"])

# --- SLIDE 29: Project Conclusion ---
add_full_slide("Conclusion", ["High accuracy models enable better supply", "Weather and Time are critical for city planning", "Scalable approach for other cities"])

# --- SLIDE 30: Thank You ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Thank You"
tf = slide.placeholders[1].text_frame
tf.text = "Thank you for the opportunity!\n\nQuestions?"

# Save
prs.save('Final_Bike_Sharing_Project_Report_30_Slides.pptx')
print("Complete 30-Slide Presentation saved.")

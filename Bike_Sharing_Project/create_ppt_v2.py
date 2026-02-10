from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Create a presentation object
prs = Presentation()

def add_full_slide(title, points, image=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    # Body text
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.word_wrap = True
    for p in points:
        p_obj = tf.add_paragraph()
        p_obj.text = p
        p_obj.level = 0
        p_obj.space_before = Pt(10)

    if image and os.path.exists(image):
        # Resize text frame to make room for image
        body_shape.width = Inches(4.5)
        slide.shapes.add_picture(image, Inches(5), Inches(1.5), height=Inches(4.5))

# Slide 1: Project Name
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Bike-sharing Rental Demand Prediction"
slide.placeholders[1].text = "Optimizing Urban Mobility with Machine Learning"

# Slide 2: Mentor & Team
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Mentor and Team Information"
tf = slide.placeholders[1].text_frame
tf.text = "Mentor: [Mentor Name]\n\nTeam Members:\n1. [Member Name 1]\n2. [Member Name 2]"

# Slide 3: Objective
add_full_slide("Project Objective", [
    "To predict the hourly demand for rental bikes based on external factors.",
    "Goal: Ensure a stable supply of bikes in urban cities.",
    "Enhance mobility comfort and reduce waiting times for commuters.",
    "Enable data-driven decision making for fleet management and pricing."
])

# Slide 4: Dataset Details
add_full_slide("Dataset Overview", [
    "Source: UCI Machine Learning Repository / Bike Sharing Dataset.",
    "Record Count: 17,379 hourly entries.",
    "Features: 17 initial attributes including temporal and environmental data.",
    "Target Variable: 'cnt' (Total number of rental bikes)."
])

# Slide 5: Data Dictionary (Part 1)
add_full_slide("Data Dictionary - Temporal Features", [
    "dteday: Date of the record.",
    "season: Season (1:springer, 2:summer, 3:fall, 4:winter).",
    "yr: Year (0: 2011, 1: 2012).",
    "mnth: Month (1 to 12).",
    "hr: Hour of the day (0 to 23).",
    "weekday: Day of the week."
])

# Slide 6: Data Dictionary (Part 2)
add_full_slide("Data Dictionary - Environmental Features", [
    "weathersit: 1:Clear, 2:Mist, 3:Light Snow, 4:Heavy Rain.",
    "temp & atemp: Normalized temperature and feeling temperature.",
    "hum: Normalized humidity levels.",
    "windspeed: Normalized wind speed.",
    "casual / registered: User type breakdown."
])

# Slide 7: Exploratory Data Analysis (EDA)
add_full_slide("Exploratory Data Analysis - Cleaning", [
    "Handled '?' symbols across multiple weather features via median imputation.",
    "Corrected date formatting issues to enable time-series extraction.",
    "Validated target consistency: verified that 'cnt' equals 'casual' + 'registered'.",
    "No duplicate records were found after data cleaning."
], "images/outliers_boxplot_cleaned.png")

# Slide 8: Data Inference - Missing Values
add_full_slide("Inference: Missing Value Treatment", [
    "Hidden missing values ('?') were detected in temperature and humidity.",
    "Imputation using the median was preferred over mean to maintain robustness against outliers.",
    "Recovered component counts for 15% of records using target sum logic."
])

# Slide 9: Data Inference - Outliers
add_full_slide("Inference: Outlier Analysis", [
    "Rental counts ('cnt') showed significant outliers during peak hours.",
    "These are not errors but naturally high demand periods (rush hours).",
    "Used robust scaling for modeling to account for these heavy-tailed distributions."
])

# Slide 10: Demand Distribution
add_full_slide("Visualizing Rental Distribution", [
    "The distribution of 'cnt' is heavily right-skewed.",
    "A large portion of hours have 0-200 rentals.",
    "Peak rentals go up to 900+ during specific conditions.",
    "Log transformation or robust models (tree-based) are ideal for this data."
], "images/dist_cnt.png")

# Slide 11: Daily Traffic Patterns
add_full_slide("Analysis: Hourly Demand Trends", [
    "Demand peaks occur at 8 AM and 5-6 PM (Commuting hours).",
    "Stable demand observed during mid-day (10 AM to 3 PM).",
    "Significant drop in rentals during late night hours (12 AM - 5 AM)."
], "images/hourly_trend.png")

# Slide 12: Relationship: Working Days
add_full_slide("Influence of Working Days & Holidays", [
    "Working days show distinctive bi-modal peaks (morning/evening rush).",
    "Holidays and weekends show a unimodal curve peaking in the afternoon.",
    "Conclusion: Operational strategies must differ between weekdays and weekends."
], "images/workingday_holiday.png")

# Slide 13: Impact: Seasonal Variation
add_full_slide("Analysis: Seasonal Demand Shift", [
    "Fall and Summer seasons witness the highest demand.",
    "Spring sees the lowest counts, likely due to weather instability.",
    "Seasonal patterns are strong drivers of demand predictability."
], "images/seasonal_demand.png")

# Slide 14: Impact: Weather Situations
add_full_slide("Analysis: Weather Impact", [
    "Clear skies (Weather 1) significantly boost rental numbers.",
    "Snow or heavy rain (Weather 3/4) lead to sharp declines in mobility.",
    "Rental systems should adjust bike availability based on live forecasts."
], "images/weather_impact.png")

# Slide 15: Environmental Analysis: Temp
add_full_slide("Temperature vs. Demand", [
    "Positive correlation observed between temperature and rentals.",
    "Comfort zone for users is between 0.4 and 0.8 normalized temperature.",
    "Extreme cold leads to almost zero casual rentals."
], "images/temp_vs_count.png")

# Slide 16: Feature Engineering - Theory
add_full_slide("Feature Engineering Strategy", [
    "Cyclic Transformations: Hour/Month converted to Sine/Cosine variables.",
    "Categorical Encoding: One-Hot Encoding for non-ordinal weather stats.",
    "Scaling: Normalized environmental variables to a 0-1 range for consistency."
])

# Slide 17: Correlation Insights
add_full_slide("Feature Correlation Analysis", [
    "Temperature has the strongest positive correlation with rentals.",
    "Humidity shows a negative correlation (higher humidity = fewer bikes).",
    "High multicollinearity detected between 'temp' and 'atemp'."
], "images/correlation_matrix_processed.png")

# Slide 18: Model Building - Approach
add_full_slide("Model Construction Workflow", [
    "Phase 1: Decision Tree Regressor (Baseline).",
    "Phase 2: Random Forest Regressor (Ensemble learning).",
    "Phase 3: Gradient Boosting Regressor (Boosting errors).",
    "Evaluation metric: RSME to penalize large prediction errors."
])

# Slide 19: Model performance Baseline
add_full_slide("Comparative Model Results", [
    "Standard Decision Tree achieved an R2 of ~0.88.",
    "Gradient Boosting showed similar performance (~0.88).",
    "Random Forest provided superior results in initial tests."
], "images/model_comparison_r2.png")

# Slide 20: Final Model Selection
add_full_slide("Selection: Random Forest Regressor", [
    "Chosen as the final model due to best RMSE and R2 scores.",
    "Random Forest handles non-linear relationships and interactions better.",
    "Effectively manages the high variance in the hourly target variable."
])

# Slide 21: Hyperparameter Tuning Impact
add_full_slide("Optimization: Hyperparameter Tuning", [
    "Conducted Randomized Search CV to find optimal tree depth and leaf sizes.",
    "Successfully balanced model complexity to avoid overfitting.",
    "The tuned model stabilized predictions for out-of-sample data."
])

# Slide 22: Model Evaluation Inferences
add_full_slide("Final Model Performance", [
    "Final R2 Score: 0.94 (Highly Accurate).",
    "Mean Absolute Error is low, indicating accurate demand estimates.",
    "The model captures both routine patterns and peak demand spikes well."
])

# Slide 23: Feature Importance Analysis
add_full_slide("Key Business Drivers", [
    "Top Predictor: The Hour of the day (hr).",
    "Second Predictor: Temperature.",
    "Third Predictor: Working Day status.",
    "Windspeed has the least effect on the count."
], "images/top_drivers.png")

# Slide 24: Prediction Accuracy
add_full_slide("Inference: Actual vs. Predicted", [
    "Scatter plot shows tight clustering around the identity line.",
    "High accuracy in mid-to-high demand forecasting.",
    "Slight variance at extremely high peaks which is expected in urban data."
], "images/actual_vs_predicted.png")

# Slide 25: Future Recommendations
add_full_slide("Business Recommendations", [
    "Deploy more bikes during 5 PM - 7 PM slots.",
    "Offer discounts during 'Weather 2' (Mist) to encourage rentals.",
    "Increase maintenance frequency during high-demand Fall months."
])

# Slide 26: Screenshot of the APP
add_full_slide("Deployment: Streamlit Web App", [
    "Developed an interactive dashboard for prediction.",
    "Allows users to input specific weather and time conditions.",
    "Generates instant demand forecasts using the optimized model."
], "images/capture_app.png") # Expected to be captured/provided

# Slide 27: Challenges Faced
add_full_slide("Challenges Faced", [
    "Data Quality: Handling hidden '?' missing values in several columns.",
    "Time Series: Converting categorical hours into meaningful numeric features.",
    "Optimization: Long training times for ensemble hyperparameter searches.",
    "Skewness: Managing the heavily right-skewed target distribution."
])

# Slide 28: Summary & Conclusion
add_full_slide("Project Conclusion", [
    "A highly accurate prediction system (94%) was achieved.",
    "Environmental and temporal factors are the primary demand drivers.",
    "Machine Learning effectively solves the urban bike-sharing supply problem."
])

# Slide 29: Thank You Note
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Thank You"
tf = slide.placeholders[1].text_frame
p = tf.add_paragraph()
p.text = "Thank you for the opportunity to work on this project!\n\nQuestions?"
p.font.size = Pt(32)
p.alignment = PP_ALIGN.CENTER

# Slide 30: Blank/Q&A
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Q & A"

# Save
prs.save('Final_Bike_Sharing_Project_Report.pptx')
print("30-Slide Presentation saved as 'Final_Bike_Sharing_Project_Report.pptx'")

from pptx import Presentation
from pptx.util import Inches, Pt
import os

# Create a presentation object
prs = Presentation()

def add_slide(title, content, image_path=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content

    if image_path and os.path.exists(image_path):
        # Adjust layout for image
        # Move text to the left
        body_shape.width = Inches(4.5)
        # Add image
        prs.slides[-1].shapes.add_picture(image_path, Inches(5), Inches(1.5), height=Inches(4))

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Bike-sharing Rental Demand Prediction"
subtitle.text = "Final Project Report\nData Science & Machine Learning Implementation"

# Slide 2: Problem Statement
add_slide("Problem Statement", 
          "- Ensure stable supply of rental bikes in urban cities.\n"
          "- Predict demand based on weather, season, holidays, and time.\n"
          "- Goal: Improve mobility comfort and customer satisfaction.")

# Slide 3: Exploratory Data Analysis
add_slide("Exploratory Data Analysis (EDA)", 
          "- Dataset: 17,379 records with 17 features.\n"
          "- Data Cleaning: Handled hidden missing values ('?') in weather features.\n"
          "- Date Parsing: Corrected DD-MM-YYYY format alignment.\n"
          "- Outlier Detection: Identified variance in counts during peak hours.")

# Slide 4: Data Visualization - Correlation
add_slide("Data Visualization", 
          "- Analyzed feature correlations using Heatmaps.\n"
          "- Observed strong relationship between temperature (temp/atemp) and rental counts.\n"
          "- Time-series trends show clear seasonal and hourly patterns.",
          "images/correlation_matrix_final.png")

# Slide 5: Feature Engineering
add_slide("Feature Engineering", 
          "- Categorical Encoding: One-Hot Encoding for 'season' and 'weather'.\n"
          "- Cyclic Encoding: Sin/Cos transformation for hour and month to capture periodicity.\n"
          "- Normalization: Min-Max Scaling for numerical features like humidity and windspeed.")

# Slide 6: Model Building & Comparison
add_slide("Model Construction", 
          "- Evaluated Decision Tree, Random Forest, and Gradient Boosting.\n"
          "- Metrics: MAE, RMSE, and R-squared.\n"
          "- Random Forest outperformed others with R2 of 0.94.",
          "images/model_comparison_r2.png")

# Slide 7: Model Optimization
add_slide("Hyperparameter Tuning", 
          "- Randomized Search performed on Random Forest.\n"
          "- Parameters tuned: n_estimators, max_depth, min_samples_split.\n"
          "- Final Model achieved high precision in capturing peak hour demand.",
          "images/feature_importance.png")

# Slide 8: Deployment & Conclusion
add_slide("Model Deployment", 
          "- Deployed via Streamlit Web Application.\n"
          "- Real-time interactive interface for demand Forecasting.\n"
          "- Stable and scalable solution for urban mobility planning.")

# Save
prs.save('Bike_Sharing_Presentation.pptx')
print("Presentation saved as 'Bike_Sharing_Presentation.pptx'")
